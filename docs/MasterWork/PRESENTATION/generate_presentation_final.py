#!/usr/bin/env python3
from __future__ import annotations

import csv
import json
from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
ROOT_DIR = BASE_DIR.parents[2]
REPORT_DIR = BASE_DIR.parent / "REPORT" / "LaTeX"
FIGURES_DIR = REPORT_DIR / "figures"
RUN_DIR = ROOT_DIR / "runs" / "detect" / "runs" / "train" / "small_detection"
VIDEO_FILE = ROOT_DIR / "videos" / "V_DRONE_009.mp4"
OUTPUT_FILE = BASE_DIR / "PRESENTATION.pptx"
RESULTS_FILE = REPORT_DIR / "acceptance_results_2026-05-13.json"
EMBLEM_FILE = BASE_DIR / "university_emblem.png"

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


BG = rgb("#F4F7FB")
INK = rgb("#101828")
MUTED = rgb("#667085")
BLUE = rgb("#2457A6")
BLUE_2 = rgb("#113B79")
CYAN = rgb("#1E9FD7")
GREEN = rgb("#287A52")
ORANGE = rgb("#B86122")
RED = rgb("#A6423D")
WHITE = rgb("#FFFFFF")
LINE = rgb("#D8E0EA")
PANEL = rgb("#FFFFFF")
PALE_BLUE = rgb("#E7F0FC")
PALE_CYAN = rgb("#E9F8FD")
PALE_GREEN = rgb("#EAF7F0")
PALE_ORANGE = rgb("#FFF3E8")

FONT_TITLE = "Arial"
FONT_BODY = "Arial"
MIN_FONT_PT = 14


def emu(value) -> int:
    return int(value)


def set_background(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def fill_shape(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def line_shape(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str = "",
    *,
    font_name: str = FONT_BODY,
    font_size: float = 14,
    bold: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
    vertical=MSO_ANCHOR.TOP,
    margins=(0.03, 0.03, 0.01, 0.01),
    auto_fit: bool = True,
):
    box = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    if auto_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = vertical
    tf.margin_left = Pt(margins[0] * 72)
    tf.margin_right = Pt(margins[1] * 72)
    tf.margin_top = Pt(margins[2] * 72)
    tf.margin_bottom = Pt(margins[3] * 72)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(max(font_size, MIN_FONT_PT))
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_line(slide, left, top, width, label: str, value: str, *, font_size=13, color=INK) -> None:
    box = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(Inches(0.24)))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.name = FONT_BODY
    run.font.size = Pt(max(font_size, MIN_FONT_PT))
    run.font.bold = True
    run.font.color.rgb = color
    run = p.add_run()
    run.text = value
    run.font.name = FONT_BODY
    run.font.size = Pt(max(font_size, MIN_FONT_PT))
    run.font.color.rgb = color


def add_paragraphs(
    textbox,
    lines: Iterable[str],
    *,
    font_size=13,
    color=INK,
    bold_first: bool = False,
    space_after=5,
    bullet: bool = False,
) -> None:
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = f"- {line}" if bullet else line
        run.font.name = FONT_BODY
        run.font.size = Pt(max(font_size, MIN_FONT_PT))
        run.font.color.rgb = color
        run.font.bold = bold_first and index == 0


def add_chip(slide, left, top, width, text: str, *, fill=PALE_BLUE, color=BLUE, font_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(width), emu(Inches(0.40)))
    fill_shape(shape, fill)
    line_shape(shape, fill, 0.5)
    shape.adjustments[0] = 0.35
    add_textbox(
        slide,
        left,
        top + Inches(0.08),
        width,
        Inches(0.20),
        text,
        font_size=font_size,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )
    return shape


def add_panel(slide, left, top, width, height, *, fill=PANEL, border=LINE, radius=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(width), emu(height))
    fill_shape(shape, fill)
    line_shape(shape, border, 0.8)
    shape.adjustments[0] = radius
    return shape


def add_header(slide, number: int, title: str, *, kicker: str | None = None) -> None:
    set_background(slide)
    if EMBLEM_FILE.exists():
        slide.shapes.add_picture(str(EMBLEM_FILE), emu(Inches(0.18)), emu(Inches(0.14)), emu(Inches(0.42)), emu(Inches(0.42)))
    add_textbox(
        slide,
        Inches(0.78),
        Inches(0.17),
        Inches(7.70),
        Inches(0.42),
        title,
        font_name=FONT_TITLE,
        font_size=20,
        bold=True,
        color=INK,
        align=PP_ALIGN.LEFT,
        vertical=MSO_ANCHOR.MIDDLE,
        margins=(0, 0, 0, 0),
    )
    if kicker:
        add_textbox(
            slide,
            Inches(0.80),
            Inches(0.55),
            Inches(7.4),
            Inches(0.24),
            kicker,
            font_size=14,
            bold=True,
            color=BLUE,
            margins=(0, 0, 0, 0),
        )
    add_textbox(
        slide,
        Inches(9.36),
        Inches(5.08),
        Inches(0.40),
        Inches(0.28),
        str(number),
        font_size=14,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        margins=(0, 0, 0, 0),
    )


def add_arrow_box(slide, left, top, width, height, text: str, *, fill=WHITE, border=BLUE, color=INK, font_size=14):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(width), emu(height))
    fill_shape(box, fill)
    line_shape(box, border, 1.0)
    box.adjustments[0] = 0.12
    add_textbox(
        slide,
        left + Inches(0.04),
        top + Inches(0.05),
        width - Inches(0.08),
        height - Inches(0.10),
        text,
        font_size=font_size,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
        margins=(0, 0, 0, 0),
    )
    return box


def add_connector(slide, x1, y1, x2, y2, *, color=BLUE, width=1.2):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(x1), emu(y1), emu(x2), emu(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_table(slide, left, top, width, height, data: list[list[str]], col_widths: list[float], *, font_size=14):
    rows = len(data)
    cols = len(data[0])
    graphic = slide.shapes.add_table(rows, cols, emu(left), emu(top), emu(width), emu(height))
    table = graphic.table
    for index, col_width in enumerate(col_widths):
        table.columns[index].width = emu(Inches(col_width))
    for row in table.rows:
        row.height = emu(height / rows)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = PALE_BLUE if r == 0 else WHITE
            cell.text_frame.clear()
            cell.text_frame.word_wrap = True
            cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            cell.text_frame.margin_left = Pt(4)
            cell.text_frame.margin_right = Pt(4)
            cell.text_frame.margin_top = Pt(2)
            cell.text_frame.margin_bottom = Pt(2)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = data[r][c]
            run.font.name = FONT_BODY
            run.font.size = Pt(max(font_size if r else font_size + 0.2, MIN_FONT_PT))
            run.font.bold = r == 0
            run.font.color.rgb = INK
    return table


def add_fit_image(slide, path: Path, left, top, width, height, *, caption: str | None = None, border=True):
    if border:
        add_panel(slide, left, top, width, height, fill=WHITE, border=LINE)
        pad = Inches(0.05)
    else:
        pad = 0
    image_top = top + pad
    image_height = height - 2 * pad
    if caption:
        image_height -= Inches(0.36)
    with Image.open(path) as img:
        iw, ih = img.size
    box_w = width - 2 * pad
    box_h = image_height
    scale = min(box_w / iw, box_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    x = left + pad + int((box_w - w) / 2)
    y = image_top + int((box_h - h) / 2)
    slide.shapes.add_picture(str(path), emu(x), emu(y), emu(w), emu(h))
    if caption:
        add_textbox(
            slide,
            left + Inches(0.03),
            top + height - Inches(0.30),
            width - Inches(0.06),
            Inches(0.24),
            caption,
            font_size=14,
            color=MUTED,
            align=PP_ALIGN.CENTER,
            margins=(0, 0, 0, 0),
        )


def add_stat(slide, left, top, width, value: str, label: str, *, color=BLUE):
    add_textbox(slide, left, top, width, Inches(0.30), value, font_size=24, bold=True, color=color, margins=(0, 0, 0, 0))
    add_textbox(slide, left, top + Inches(0.40), width, Inches(0.24), label, font_size=14, color=MUTED, margins=(0, 0, 0, 0))


def load_acceptance() -> dict:
    return json.loads(RESULTS_FILE.read_text(encoding="utf-8"))


def load_yolo_metrics() -> dict[str, float]:
    with (RUN_DIR / "results.csv").open(encoding="utf-8") as file:
        rows = list(csv.DictReader(file))
    last = rows[-1]
    best = max(rows, key=lambda row: float(row["metrics/mAP50(B)"]))
    return {
        "epochs": float(last["epoch"]),
        "precision": float(last["metrics/precision(B)"]),
        "recall": float(last["metrics/recall(B)"]),
        "map50": float(last["metrics/mAP50(B)"]),
        "map5095": float(last["metrics/mAP50-95(B)"]),
        "best_epoch": float(best["epoch"]),
        "best_map50": float(best["metrics/mAP50(B)"]),
    }


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    if EMBLEM_FILE.exists():
        slide.shapes.add_picture(str(EMBLEM_FILE), emu(Inches(0.35)), emu(Inches(0.25)), emu(Inches(1.10)), emu(Inches(1.10)))
    add_textbox(slide, Inches(1.65), Inches(0.30), Inches(7.85), Inches(0.28), "Московский авиационный институт", font_size=14, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(1.65), Inches(0.64), Inches(7.85), Inches(0.28), "Институт №8, кафедра №806", font_size=14, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(0.90), Inches(1.58), Inches(8.15), Inches(0.24), "Выпускная квалификационная работа на тему:", font_size=17, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    title = (
        "«Разработка и экспериментальная оценка платформы\n"
        "видеоаналитики реального времени для обнаружения,\n"
        "сопровождения и событийного анализа объектов\n"
        "в системах охранного видеонаблюдения»"
    )
    add_textbox(slide, Inches(0.85), Inches(1.95), Inches(8.30), Inches(1.32), title, font_size=21, bold=True, color=BLUE, align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE, margins=(0, 0, 0, 0))
    add_rich_line(slide, Inches(0.78), Inches(4.12), Inches(8.30), "Студент: ", "Нурмухамедов А.Б., М8О-214СВ-24", font_size=14)
    add_rich_line(slide, Inches(0.78), Inches(4.44), Inches(8.30), "Руководитель: ", "Гаврилов Константин Юрьевич, д.т.н., профессор", font_size=14)
    add_textbox(slide, Inches(4.15), Inches(5.04), Inches(1.70), Inches(0.26), "Москва, 2026", font_size=14, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def slide_relevance(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 2, "Актуальность")
    add_panel(slide, Inches(0.70), Inches(1.05), Inches(2.55), Inches(3.35), fill=PALE_BLUE, border=PALE_BLUE)
    add_textbox(slide, Inches(0.98), Inches(1.32), Inches(2.00), Inches(0.62), "Камер больше", font_size=24, bold=True, color=BLUE_2, align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE)
    for i, y in enumerate([2.25, 2.88, 3.51], start=1):
        add_arrow_box(slide, Inches(1.08), Inches(y), Inches(0.95), Inches(0.38), str(i), fill=WHITE, border=BLUE)
        add_connector(slide, Inches(2.05), Inches(y + 0.19), Inches(2.45), Inches(3.06))
    add_arrow_box(slide, Inches(2.26), Inches(2.78), Inches(0.82), Inches(0.56), "Опер.", fill=WHITE, border=BLUE)

    add_panel(slide, Inches(3.65), Inches(1.05), Inches(2.35), Inches(3.35), fill=WHITE)
    add_textbox(slide, Inches(3.95), Inches(1.28), Inches(1.75), Inches(0.30), "Риски", font_size=18, bold=True, color=RED, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_stat(slide, Inches(4.02), Inches(1.85), Inches(1.60), "latency", "задержка", color=RED)
    add_stat(slide, Inches(4.02), Inches(2.75), Inches(1.60), "miss", "пропуск", color=ORANGE)
    add_stat(slide, Inches(4.02), Inches(3.65), Inches(1.60), "load", "нагрузка", color=BLUE)

    add_panel(slide, Inches(6.40), Inches(1.05), Inches(2.85), Inches(3.35), fill=WHITE)
    add_textbox(slide, Inches(6.70), Inches(1.28), Inches(2.25), Inches(0.30), "Нужен контур", font_size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    for idx, (label, color) in enumerate([("detect", BLUE), ("track", CYAN), ("rules", GREEN), ("notify", ORANGE)]):
        add_chip(slide, Inches(7.08), Inches(1.90 + idx * 0.58), Inches(1.45), label, fill=PALE_CYAN, color=color)
    add_textbox(slide, Inches(0.95), Inches(4.82), Inches(8.10), Inches(0.28), "Объект -> трек -> правило -> событие -> уведомление", font_size=16, bold=True, color=BLUE_2, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def slide_goal_tasks(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 3, "Цель и задачи")
    add_panel(slide, Inches(0.58), Inches(1.02), Inches(3.78), Inches(3.92), fill=WHITE)
    add_textbox(slide, Inches(0.82), Inches(1.28), Inches(3.30), Inches(0.30), "Цель", font_size=18, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(0.82), Inches(1.88), Inches(3.30), Inches(1.20), "Real-time мониторинг охраняемых зон", font_size=26, bold=True, color=INK, vertical=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.82), Inches(3.55), Inches(3.25), Inches(0.55), "Объект: видеопотоки\nПредмет: CV + архитектура", font_size=16, color=BLUE_2, vertical=MSO_ANCHOR.MIDDLE, margins=(0, 0, 0, 0))

    tasks = [
        ("1", "треб."),
        ("2", "CV"),
        ("3", "арх."),
        ("4", "код"),
        ("5", "тест"),
        ("6", "план"),
    ]
    lefts = [4.88, 6.35, 7.82]
    tops = [1.25, 2.90]
    for idx, (num, label) in enumerate(tasks):
        col = idx % 3
        row = idx // 3
        add_panel(slide, Inches(lefts[col]), Inches(tops[row]), Inches(1.15), Inches(1.05), fill=[PALE_BLUE, PALE_CYAN, PALE_GREEN, PALE_ORANGE, WHITE, PALE_BLUE][idx])
        add_textbox(slide, Inches(lefts[col] + 0.10), Inches(tops[row] + 0.12), Inches(0.32), Inches(0.20), num, font_size=16, bold=True, color=BLUE, margins=(0, 0, 0, 0))
        add_textbox(slide, Inches(lefts[col] + 0.08), Inches(tops[row] + 0.52), Inches(0.98), Inches(0.28), label, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def slide_expected_result(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 4, "Ключевые решения")
    data = [
        ["Слой", "Решение", "Зачем"],
        ["CV", "YOLO + DeepSORT", "детекция и треки"],
        ["Rules", "zone / line / dwell", "интерпретируемые события"],
        ["Delivery", "outbox + retry", "сохранность событий"],
        ["Access", "JWT + RBAC", "роли оператора"],
        ["Ops", "health + metrics", "наблюдаемость"],
    ]
    add_table(slide, Inches(0.58), Inches(1.10), Inches(8.85), Inches(2.55), data, [2.10, 3.20, 3.55], font_size=14)
    add_panel(slide, Inches(0.58), Inches(3.82), Inches(8.85), Inches(0.72), fill=PALE_BLUE, border=PALE_BLUE)
    add_textbox(slide, Inches(0.88), Inches(3.98), Inches(8.20), Inches(0.28), "Новизна: real-time CV + надежная событийная доставка.", font_size=18, bold=True, color=BLUE_2, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_chip(slide, Inches(1.18), Inches(4.78), Inches(1.70), "capture")
    add_chip(slide, Inches(3.18), Inches(4.78), Inches(1.70), "track")
    add_chip(slide, Inches(5.18), Inches(4.78), Inches(1.70), "rules")
    add_chip(slide, Inches(7.18), Inches(4.78), Inches(1.70), "notify")


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 5, "Архитектура")
    add_textbox(slide, Inches(0.68), Inches(0.90), Inches(1.50), Inches(0.26), "Control", font_size=14, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(0.68), Inches(2.32), Inches(1.50), Inches(0.26), "Data", font_size=14, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(0.68), Inches(3.98), Inches(1.50), Inches(0.26), "Integration", font_size=14, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    add_arrow_box(slide, Inches(0.78), Inches(1.15), Inches(1.05), Inches(0.50), "UI", fill=WHITE)
    add_arrow_box(slide, Inches(2.25), Inches(1.15), Inches(1.35), Inches(0.50), "FastAPI\nJWT/RBAC", fill=WHITE, font_size=10.5)
    add_arrow_box(slide, Inches(4.10), Inches(1.15), Inches(1.50), Inches(0.50), "PostgreSQL", fill=WHITE)
    add_arrow_box(slide, Inches(6.52), Inches(1.15), Inches(2.12), Inches(0.50), "Prometheus / Grafana", fill=WHITE, font_size=10.5)
    add_connector(slide, Inches(1.83), Inches(1.40), Inches(2.25), Inches(1.40))
    add_connector(slide, Inches(3.60), Inches(1.40), Inches(4.10), Inches(1.40))
    add_connector(slide, Inches(5.60), Inches(1.40), Inches(6.52), Inches(1.40))

    add_arrow_box(slide, Inches(0.78), Inches(2.65), Inches(1.12), Inches(0.55), "RTSP /\nMP4", fill=PALE_BLUE, font_size=10.5)
    add_arrow_box(slide, Inches(2.25), Inches(2.65), Inches(1.45), Inches(0.55), "analytics-\nworker", fill=PALE_CYAN, font_size=10.5)
    add_arrow_box(slide, Inches(4.15), Inches(2.65), Inches(1.35), Inches(0.55), "events +\noutbox", fill=WHITE, font_size=10.5)
    add_arrow_box(slide, Inches(6.05), Inches(2.65), Inches(1.32), Inches(0.55), "RabbitMQ", fill=PALE_BLUE, font_size=10.5)
    add_arrow_box(slide, Inches(7.85), Inches(2.65), Inches(1.35), Inches(0.55), "integration-\nworker", fill=WHITE, font_size=10.2)
    for x1, x2 in [(1.90, 2.25), (3.70, 4.15), (5.50, 6.05), (7.37, 7.85)]:
        add_connector(slide, Inches(x1), Inches(2.92), Inches(x2), Inches(2.92))
    add_connector(slide, Inches(4.80), Inches(1.65), Inches(4.80), Inches(2.65))

    add_arrow_box(slide, Inches(5.92), Inches(4.12), Inches(1.62), Inches(0.55), "webhook", fill=WHITE)
    add_arrow_box(slide, Inches(7.92), Inches(4.12), Inches(1.25), Inches(0.55), "DLQ", fill=PALE_ORANGE, border=ORANGE, color=ORANGE)
    add_connector(slide, Inches(8.52), Inches(3.20), Inches(6.75), Inches(4.12))
    add_connector(slide, Inches(8.52), Inches(3.20), Inches(8.52), Inches(4.12), color=ORANGE)


def slide_pipeline(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 6, "Конвейер обработки")
    steps = [
        ("capture", ""),
        ("latest", ""),
        ("YOLO", ""),
        ("DeepSORT", ""),
        ("rules", ""),
        ("outbox", ""),
        ("webhook", ""),
    ]
    x = 0.50
    for idx, (title, sub) in enumerate(steps):
        add_arrow_box(slide, Inches(x), Inches(1.18), Inches(1.02), Inches(0.62), title, fill=WHITE if idx not in [1, 4] else PALE_BLUE, font_size=14)
        if idx < len(steps) - 1:
            add_connector(slide, Inches(x + 1.02), Inches(1.49), Inches(x + 1.22), Inches(1.49))
        x += 1.28
    add_panel(slide, Inches(0.78), Inches(2.45), Inches(8.50), Inches(1.25), fill=PALE_GREEN, border=PALE_GREEN)
    add_textbox(slide, Inches(1.10), Inches(2.78), Inches(7.86), Inches(0.46), "Свежий кадр важнее старой очереди", font_size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE)
    add_chip(slide, Inches(1.45), Inches(4.42), Inches(1.55), "low latency")
    add_chip(slide, Inches(3.45), Inches(4.42), Inches(1.55), "track_id")
    add_chip(slide, Inches(5.45), Inches(4.42), Inches(1.55), "dedup_key")
    add_chip(slide, Inches(7.45), Inches(4.42), Inches(1.20), "retry")


def slide_rules(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 7, "Правила событий")
    add_panel(slide, Inches(0.58), Inches(1.05), Inches(2.55), Inches(3.40), fill=WHITE)
    add_textbox(slide, Inches(0.82), Inches(1.26), Inches(2.00), Inches(0.20), "zone_enter", font_size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    zone = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, emu(Inches(1.10)), emu(Inches(1.82)), emu(Inches(1.50)), emu(Inches(1.15)))
    fill_shape(zone, PALE_BLUE)
    line_shape(zone, BLUE, 1.2)
    add_arrow_box(slide, Inches(1.40), Inches(3.25), Inches(0.70), Inches(0.38), "track", fill=PALE_CYAN, font_size=9.5)
    add_textbox(slide, Inches(0.86), Inches(3.85), Inches(1.95), Inches(0.28), "Вход", font_size=16, color=MUTED, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))

    add_panel(slide, Inches(3.35), Inches(1.05), Inches(2.55), Inches(3.40), fill=WHITE)
    add_textbox(slide, Inches(3.60), Inches(1.26), Inches(2.00), Inches(0.20), "line_cross", font_size=15, bold=True, color=CYAN, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_connector(slide, Inches(3.78), Inches(2.55), Inches(5.45), Inches(2.55), color=CYAN, width=2.4)
    add_connector(slide, Inches(4.08), Inches(3.18), Inches(5.12), Inches(1.88), color=BLUE, width=1.8)
    add_textbox(slide, Inches(3.65), Inches(3.85), Inches(2.00), Inches(0.28), "Рубеж", font_size=16, color=MUTED, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))

    add_panel(slide, Inches(6.12), Inches(1.05), Inches(2.55), Inches(3.40), fill=WHITE)
    add_textbox(slide, Inches(6.36), Inches(1.26), Inches(2.00), Inches(0.20), "dwell_time", font_size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu(Inches(6.96)), emu(Inches(1.82)), emu(Inches(1.05)), emu(Inches(1.05)))
    fill_shape(circle, PALE_GREEN)
    line_shape(circle, GREEN, 1.5)
    add_textbox(slide, Inches(7.13), Inches(2.20), Inches(0.70), Inches(0.15), "T > N", font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(6.40), Inches(3.85), Inches(2.00), Inches(0.28), "Задержка", font_size=16, color=MUTED, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))

    add_panel(slide, Inches(0.85), Inches(4.78), Inches(7.80), Inches(0.40), fill=PALE_BLUE, border=PALE_BLUE)
    add_textbox(slide, Inches(1.00), Inches(4.86), Inches(7.50), Inches(0.22), "dedup_key подавляет повторы", font_size=16, bold=True, color=BLUE_2, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def slide_reliability(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 8, "Доставка событий")
    add_arrow_box(slide, Inches(0.68), Inches(1.22), Inches(1.55), Inches(0.56), "Rule Engine", fill=PALE_CYAN)
    add_arrow_box(slide, Inches(2.82), Inches(1.22), Inches(1.55), Inches(0.56), "events", fill=WHITE)
    add_arrow_box(slide, Inches(4.92), Inches(1.22), Inches(1.60), Inches(0.56), "event_outbox", fill=PALE_BLUE)
    add_arrow_box(slide, Inches(7.05), Inches(1.22), Inches(1.55), Inches(0.56), "RabbitMQ", fill=WHITE)
    for x1, x2 in [(2.23, 2.82), (4.37, 4.92), (6.52, 7.05)]:
        add_connector(slide, Inches(x1), Inches(1.50), Inches(x2), Inches(1.50))
    add_textbox(slide, Inches(3.25), Inches(1.92), Inches(2.80), Inches(0.28), "одна транзакция", font_size=16, color=GREEN, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))

    add_panel(slide, Inches(0.70), Inches(2.65), Inches(3.85), Inches(1.55), fill=WHITE)
    add_textbox(slide, Inches(0.95), Inches(2.86), Inches(3.35), Inches(0.28), "Состояния outbox", font_size=16, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    for i, state in enumerate(["new", "retry", "ok", "dlq"]):
        add_chip(slide, Inches(0.95 + i * 0.86), Inches(3.34), Inches(0.78), state, fill=PALE_ORANGE if state in ["retry", "dlq"] else PALE_BLUE, color=ORANGE if state in ["retry", "dlq"] else BLUE, font_size=14)

    add_panel(slide, Inches(4.90), Inches(2.65), Inches(4.10), Inches(1.55), fill=WHITE)
    add_textbox(slide, Inches(5.15), Inches(2.86), Inches(3.50), Inches(0.28), "Проверено", font_size=16, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    checks = ["RabbitMQ", "webhook", "retry", "DLQ=0"]
    for i, check in enumerate(checks):
        add_textbox(slide, Inches(5.15 + (i % 2) * 1.85), Inches(3.28 + (i // 2) * 0.42), Inches(1.65), Inches(0.26), check, font_size=15, color=INK, margins=(0, 0, 0, 0))
    add_panel(slide, Inches(0.92), Inches(4.70), Inches(8.00), Inches(0.38), fill=PALE_GREEN, border=PALE_GREEN)
    add_textbox(slide, Inches(1.08), Inches(4.77), Inches(7.68), Inches(0.22), "Сбой != потеря события", font_size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def slide_ui(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 8, "UI и демо")
    add_fit_image(slide, FIGURES_DIR / "ui-dashboard.png", Inches(0.42), Inches(1.00), Inches(2.18), Inches(1.58), caption="Dashboard")
    add_fit_image(slide, FIGURES_DIR / "ui-demo.png", Inches(2.78), Inches(1.00), Inches(2.18), Inches(1.58), caption="Demo")
    add_fit_image(slide, FIGURES_DIR / "ui-events.png", Inches(5.14), Inches(1.00), Inches(2.18), Inches(1.58), caption="Events")
    add_fit_image(slide, FIGURES_DIR / "ui-status.png", Inches(7.50), Inches(1.00), Inches(2.18), Inches(1.58), caption="Status")
    features = [
        ("Acknowledge", ""),
        ("Evidence", ""),
        ("Routes", ""),
        ("Status", ""),
    ]
    x = 1.22
    for title, text_value in features:
        add_panel(slide, Inches(x), Inches(3.35), Inches(1.55), Inches(0.88), fill=WHITE)
        add_textbox(slide, Inches(x + 0.10), Inches(3.62), Inches(1.35), Inches(0.24), title, font_size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
        x += 1.88
    add_textbox(slide, Inches(0.78), Inches(4.82), Inches(8.30), Inches(0.24), "PATCH ack + GET evidence работают через API", font_size=18, bold=True, color=BLUE_2, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def slide_stack(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 7, "Реализация")
    layers = [
        ("API", "FastAPI\nRBAC"),
        ("DB", "PostgreSQL\nAlembic"),
        ("Workers", "analytics\nintegration"),
        ("UI", "React\nTypeScript"),
        ("Ops", "Docker\nPrometheus"),
    ]
    colors = [PALE_BLUE, PALE_CYAN, PALE_GREEN, PALE_ORANGE, WHITE]
    x = 0.58
    for idx, (title, body) in enumerate(layers):
        add_panel(slide, Inches(x), Inches(1.20), Inches(1.65), Inches(2.35), fill=colors[idx])
        add_textbox(slide, Inches(x + 0.16), Inches(1.48), Inches(1.30), Inches(0.22), title, font_size=16, bold=True, color=BLUE_2, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
        add_textbox(slide, Inches(x + 0.16), Inches(2.20), Inches(1.30), Inches(0.70), body, font_size=15, color=INK, align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE, margins=(0, 0, 0, 0))
        x += 1.83
    add_panel(slide, Inches(1.65), Inches(4.18), Inches(6.65), Inches(0.52), fill=WHITE)
    add_textbox(slide, Inches(1.85), Inches(4.32), Inches(6.25), Inches(0.22), "make verify  |  make acceptance-e2e", font_size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def slide_training_setup(prs: Presentation, metrics: dict[str, float]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 11, "Обучение детектора")
    add_fit_image(slide, FIGURES_DIR / "yolo_training_results.png", Inches(0.55), Inches(1.05), Inches(5.35), Inches(3.55), caption="Динамика обучения YOLO")
    add_panel(slide, Inches(6.20), Inches(1.05), Inches(3.15), Inches(3.55), fill=WHITE)
    add_textbox(slide, Inches(6.45), Inches(1.25), Inches(2.65), Inches(0.28), "Профиль", font_size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    lines = [
        "Small Detection",
        "bird + drone",
        "100 эпох",
    ]
    add_paragraphs(slide.shapes.add_textbox(emu(Inches(6.58)), emu(Inches(1.78)), emu(Inches(2.40)), emu(Inches(1.20))), lines, font_size=16, color=INK, space_after=6, bullet=True)
    add_stat(slide, Inches(6.55), Inches(3.45), Inches(1.20), f"{metrics['map50']:.3f}", "mAP50", color=GREEN)
    add_stat(slide, Inches(7.82), Inches(3.45), Inches(1.20), f"{metrics['best_map50']:.3f}", "best", color=BLUE)


def slide_validation(prs: Presentation, metrics: dict[str, float]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 12, "Валидация детектора")
    add_fit_image(slide, FIGURES_DIR / "yolo_confusion_matrix.png", Inches(0.60), Inches(1.02), Inches(3.08), Inches(3.78), caption="Матрица ошибок")
    add_fit_image(slide, FIGURES_DIR / "yolo_validation_predictions.jpg", Inches(3.95), Inches(1.02), Inches(3.08), Inches(3.78), caption="Валидационные предсказания")
    add_panel(slide, Inches(7.28), Inches(1.02), Inches(2.08), Inches(3.78), fill=WHITE)
    add_stat(slide, Inches(7.55), Inches(1.43), Inches(1.35), f"{metrics['precision']:.3f}", "precision", color=BLUE)
    add_stat(slide, Inches(7.55), Inches(2.30), Inches(1.35), f"{metrics['recall']:.3f}", "recall", color=CYAN)
    add_stat(slide, Inches(7.55), Inches(3.17), Inches(1.35), f"{metrics['map5095']:.3f}", "mAP50-95", color=GREEN)


def slide_video_demo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 13, "Видео-демонстрация")
    poster = FIGURES_DIR / "drone_demo_detection.png"
    add_panel(slide, Inches(0.55), Inches(1.02), Inches(5.45), Inches(3.85), fill=WHITE)
    slide.shapes.add_movie(
        str(VIDEO_FILE),
        emu(Inches(0.77)),
        emu(Inches(1.20)),
        emu(Inches(5.00)),
        emu(Inches(3.35)),
        poster_frame_image=str(poster),
        mime_type="video/mp4",
    )
    add_panel(slide, Inches(6.30), Inches(1.02), Inches(3.05), Inches(3.85), fill=WHITE)
    add_textbox(slide, Inches(6.55), Inches(1.25), Inches(2.55), Inches(0.28), "MP4 внутри PPTX", font_size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_chip(slide, Inches(6.85), Inches(1.95), Inches(1.85), "drone")
    add_chip(slide, Inches(6.85), Inches(2.55), Inches(1.85), "best.pt")
    add_chip(slide, Inches(6.85), Inches(3.15), Inches(1.85), "zone_enter")
    add_stat(slide, Inches(6.65), Inches(3.45), Inches(1.18), "10.07с", "длительность", color=BLUE)
    add_stat(slide, Inches(7.92), Inches(3.45), Inches(1.05), "30 fps", "частота", color=CYAN)


def slide_methodology(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 9, "Методика апробации")
    data = [
        ["Контур", "Проверка", "Критерий"],
        ["CPU/E2E", "event -> outbox -> webhook", "PASS"],
        ["Reliability", "RabbitMQ / webhook outage", "restore"],
        ["Security", "JWT / RBAC viewer", "403"],
        ["Demo", "БПЛА + UI", "событие"],
    ]
    add_table(slide, Inches(0.60), Inches(1.03), Inches(8.80), Inches(1.95), data, [1.85, 4.10, 2.85], font_size=14)
    add_panel(slide, Inches(0.80), Inches(3.35), Inches(2.25), Inches(1.10), fill=PALE_BLUE, border=PALE_BLUE)
    add_textbox(slide, Inches(0.98), Inches(3.55), Inches(1.85), Inches(0.18), "Latency", font_size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(1.02), Inches(3.90), Inches(1.78), Inches(0.24), "functional", font_size=15, color=INK, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_panel(slide, Inches(3.42), Inches(3.35), Inches(2.25), Inches(1.10), fill=PALE_GREEN, border=PALE_GREEN)
    add_textbox(slide, Inches(3.60), Inches(3.55), Inches(1.85), Inches(0.18), "Delivery", font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(3.66), Inches(3.90), Inches(1.72), Inches(0.24), "success / retry", font_size=15, color=INK, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_panel(slide, Inches(6.05), Inches(3.35), Inches(2.25), Inches(1.10), fill=PALE_ORANGE, border=PALE_ORANGE)
    add_textbox(slide, Inches(6.23), Inches(3.55), Inches(1.85), Inches(0.18), "Security", font_size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(6.28), Inches(3.90), Inches(1.75), Inches(0.24), "JWT / RBAC", font_size=15, color=INK, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def slide_results(prs: Presentation, results: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 10, "CPU acceptance")
    add_stat(slide, Inches(0.72), Inches(0.95), Inches(1.25), "PASS", "общий статус", color=GREEN)
    add_stat(slide, Inches(2.05), Inches(0.95), Inches(1.48), f"{results['duration_sec']:.3f}с", "длительность", color=BLUE)
    add_stat(slide, Inches(3.58), Inches(0.95), Inches(1.25), "403", "viewer write", color=ORANGE)
    chart_data = CategoryChartData()
    chart_data.categories = ["functional", "webhook", "RabbitMQ"]
    chart_data.add_series(
        "seconds",
        (
            results["timings_sec"]["functional_delivery"],
            results["timings_sec"]["webhook_recovery"],
            results["timings_sec"]["rabbitmq_recovery"],
        ),
    )
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, emu(Inches(0.72)), emu(Inches(2.05)), emu(Inches(4.30)), emu(Inches(2.55)), chart_data)
    chart = graphic.chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(14)
    chart.value_axis.tick_labels.font.size = Pt(14)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.number_format = '0.000 "c"'
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Время доставки / восстановления"
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
    data = [
        ["Проверка", "Факт"],
        ["E2E", "published"],
        ["RabbitMQ", "retry -> published"],
        ["Webhook", "failed -> success"],
        ["DLQ", "0"],
    ]
    add_table(slide, Inches(5.28), Inches(1.30), Inches(4.08), Inches(2.92), data, [1.35, 2.73], font_size=14)
    add_textbox(slide, Inches(5.35), Inches(4.50), Inches(3.92), Inches(0.28), "Подтверждено: E2E + reliability + RBAC", font_size=16, color=GREEN, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def slide_artifacts(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 16, "Артефакты")
    data = [
        ["Артефакт", "Статус"],
        ["PPTX", "готов"],
        ["Acceptance JSON", "PASS"],
        ["Демо-видео", "встроено"],
        ["Модель", "best.pt"],
        ["Превью", "11 PNG"],
        ["Речь", "обновлена"],
    ]
    add_table(slide, Inches(0.68), Inches(1.05), Inches(8.65), Inches(3.00), data, [3.20, 5.45], font_size=14)
    add_panel(slide, Inches(0.95), Inches(4.45), Inches(7.95), Inches(0.50), fill=PALE_BLUE, border=PALE_BLUE)
    add_textbox(slide, Inches(1.12), Inches(4.58), Inches(7.60), Inches(0.24), "Проверяемые файлы лежат в проекте.", font_size=18, bold=True, color=BLUE_2, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def slide_conclusion(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, 11, "Выводы и развитие")
    add_panel(slide, Inches(0.70), Inches(1.00), Inches(3.80), Inches(3.25), fill=WHITE)
    add_textbox(slide, Inches(0.98), Inches(1.22), Inches(3.20), Inches(0.20), "Получено", font_size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_paragraphs(
        slide.shapes.add_textbox(emu(Inches(0.98)), emu(Inches(1.62)), emu(Inches(3.20)), emu(Inches(1.85))),
        [
            "real-time",
            "events",
            "outbox",
            "БПЛА",
        ],
        font_size=16,
        color=INK,
        space_after=6,
        bullet=True,
    )
    add_panel(slide, Inches(5.05), Inches(1.00), Inches(3.80), Inches(3.25), fill=WHITE)
    add_textbox(slide, Inches(5.35), Inches(1.22), Inches(3.20), Inches(0.20), "Дальнейшее развитие", font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))
    add_paragraphs(
        slide.shapes.add_textbox(emu(Inches(5.35)), emu(Inches(1.62)), emu(Inches(3.20)), emu(Inches(1.85))),
        [
            "данные",
            "ReID",
            "storage",
            "ONNX/TensorRT",
        ],
        font_size=16,
        color=INK,
        space_after=6,
        bullet=True,
    )
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(Inches(0.70)), emu(Inches(4.72)), emu(Inches(8.15)), emu(Inches(0.42)))
    fill_shape(band, BLUE)
    line_shape(band, BLUE, 0.5)
    add_textbox(slide, Inches(0.70), Inches(4.80), Inches(8.15), Inches(0.22), "Цель работы достигнута.", font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, margins=(0, 0, 0, 0))


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    results = load_acceptance()
    slide_title(prs)
    slide_relevance(prs)
    slide_goal_tasks(prs)
    slide_expected_result(prs)
    slide_architecture(prs)
    slide_pipeline(prs)
    slide_stack(prs)
    slide_ui(prs)
    slide_methodology(prs)
    slide_results(prs, results)
    slide_conclusion(prs)
    return prs


if __name__ == "__main__":
    presentation = build_presentation()
    presentation.save(OUTPUT_FILE)
    print(OUTPUT_FILE)
