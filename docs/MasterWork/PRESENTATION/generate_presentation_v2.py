#!/usr/bin/env python3
from __future__ import annotations

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
REPORT_DIR = BASE_DIR.parent / "REPORT" / "LaTeX"
OUTPUT_FILE = BASE_DIR / "PRESENTATION.pptx"
METADATA_FILE = REPORT_DIR / "metadata.tex"
RESULTS_FILE = REPORT_DIR / "acceptance_results_2026-05-13.json"
ASSETS_DIR = BASE_DIR / "figma_assets"
EMBLEM_FILE = BASE_DIR / "university_emblem.png"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


BG = rgb("#F3F3F3")
BLUE = rgb("#2B59A2")
LIGHT_BLUE = rgb("#E9F0FB")
LIGHTER_BLUE = rgb("#F5F8FE")
BLACK = rgb("#111111")
GRAY = rgb("#5E5E5E")
LIGHT_GRAY = rgb("#D9DDE5")
WHITE = rgb("#FFFFFF")
GREEN = rgb("#3F7D4E")
ORANGE = rgb("#B46B22")
RED = rgb("#9F3D39")

FONT_TITLE = "Oswald"
FONT_BODY = "Arial"

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)


def emu(value) -> int:
    return int(value)


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def fill_shape(shape, color: RGBColor) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def line_shape(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    *,
    font_name=FONT_BODY,
    font_size=14,
    bold=False,
    color=BLACK,
    align=PP_ALIGN.LEFT,
    vertical=MSO_ANCHOR.TOP,
    margins=(0.03, 0.03, 0.01, 0.01),
):
    box = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = vertical
    tf.margin_left = Pt(margins[0] * 72)
    tf.margin_right = Pt(margins[1] * 72)
    tf.margin_top = Pt(margins[2] * 72)
    tf.margin_bottom = Pt(margins[3] * 72)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(
    textbox,
    lines: list[str],
    *,
    font_size=13,
    color=BLACK,
    level=0,
    bullet_char="•",
    space_after=7,
) -> None:
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = level
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = f"{bullet_char} {line}"
        run.font.name = FONT_BODY
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def add_paragraph_lines(
    textbox,
    lines: list[str],
    *,
    font_size=13,
    color=BLACK,
    bold_first=False,
    align=PP_ALIGN.LEFT,
    space_after=5,
) -> None:
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.name = FONT_BODY
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        if bold_first and index == 0:
            run.font.bold = True


def add_label_value_line(slide, left, top, width, label: str, value: str, *, font_size=13, color=BLACK):
    box = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(Inches(0.30)))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run_label = p.add_run()
    run_label.text = label
    run_label.font.name = FONT_BODY
    run_label.font.size = Pt(font_size)
    run_label.font.bold = True
    run_label.font.color.rgb = color

    run_value = p.add_run()
    run_value.text = value
    run_value.font.name = FONT_BODY
    run_value.font.size = Pt(font_size)
    run_value.font.color.rgb = color
    return box


def add_reference_badge(slide, *, large=False) -> None:
    size = Inches(0.90 if large else 0.78)
    left = Inches(0.06)
    top = Inches(0.06)
    if EMBLEM_FILE.exists():
        slide.shapes.add_picture(str(EMBLEM_FILE), emu(left), emu(top), emu(size), emu(size))
        return

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(size), emu(size))
    fill_shape(badge, BLUE)
    line_shape(badge, BLUE, 0.8)
    badge.adjustments[0] = 0.18


def add_slide_chrome(slide, number: int, title: str) -> None:
    set_background(slide)
    add_reference_badge(slide)
    add_textbox(
        slide,
        Inches(0.34),
        Inches(0.44),
        Inches(9.0),
        Inches(0.34),
        title,
        font_name=FONT_TITLE,
        font_size=20,
        bold=True,
        color=BLACK,
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
        margins=(0, 0, 0, 0),
    )
    add_textbox(
        slide,
        Inches(9.28),
        Inches(5.08),
        Inches(0.35),
        Inches(0.20),
        str(number),
        font_name=FONT_BODY,
        font_size=12,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
        margins=(0, 0, 0, 0),
    )


def add_note(slide, left, top, width, text: str, *, height=0.55, font_size=11) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(width), emu(Inches(height)))
    fill_shape(box, WHITE)
    line_shape(box, LIGHT_GRAY, 0.8)
    box.adjustments[0] = 0.05
    add_textbox(
        slide,
        left + Inches(0.10),
        top + Inches(0.08),
        width - Inches(0.20),
        Inches(height - 0.16),
        text,
        font_name=FONT_BODY,
        font_size=font_size,
        color=GRAY,
    )


def add_content_box(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(left), emu(top), emu(width), emu(height))
    fill_shape(shape, WHITE)
    line_shape(shape, LIGHT_GRAY, 0.8)
    return shape


def add_connector(slide, x1, y1, x2, y2, *, color=BLUE, width=1.2):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(x1), emu(y1), emu(x2), emu(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_arrow_box(slide, left, top, width, height, text, *, fill=LIGHT_BLUE, border=BLUE, font_size=12):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(left), emu(top), emu(width), emu(height))
    fill_shape(box, fill)
    line_shape(box, border, 0.9)
    box.adjustments[0] = 0.10
    add_textbox(
        slide,
        left + Inches(0.04),
        top + Inches(0.12),
        width - Inches(0.08),
        height - Inches(0.16),
        text,
        font_name=FONT_BODY,
        font_size=font_size,
        color=BLACK,
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
    )
    return box


def add_table(slide, left, top, width, height, data: list[list[str]], col_widths: list[float], *, header_fill=LIGHT_BLUE):
    rows = len(data)
    cols = len(data[0])
    graphic = slide.shapes.add_table(rows, cols, emu(left), emu(top), emu(width), emu(height))
    table = graphic.table

    for index, col_width in enumerate(col_widths):
        table.columns[index].width = emu(Inches(col_width))

    row_height = height / rows
    for row in table.rows:
        row.height = emu(row_height)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else WHITE
            cell.text_frame.clear()
            cell.text_frame.word_wrap = True
            cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            cell.text_frame.margin_left = Pt(4)
            cell.text_frame.margin_right = Pt(4)
            cell.text_frame.margin_top = Pt(3)
            cell.text_frame.margin_bottom = Pt(3)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = data[r][c]
            run.font.name = FONT_BODY
            run.font.size = Pt(11 if r == 0 else 10.5)
            run.font.color.rgb = BLACK
            run.font.bold = r == 0
    return table


def add_image(slide, path: Path, left, top, width, height, caption: str):
    frame = add_content_box(slide, left, top, width, height)
    slide.shapes.add_picture(
        str(path),
        emu(left + Inches(0.03)),
        emu(top + Inches(0.03)),
        emu(width - Inches(0.06)),
        emu(height - Inches(0.24)),
    )
    add_textbox(
        slide,
        left,
        top + height - Inches(0.18),
        width,
        Inches(0.16),
        caption,
        font_name=FONT_BODY,
        font_size=10,
        color=GRAY,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )
    return frame


def add_horizontal_bar_chart(slide, left, top, width, height, items: list[tuple[str, float, RGBColor]], *, max_value: float, title: str):
    add_textbox(
        slide,
        left,
        top,
        width,
        Inches(0.24),
        title,
        font_name=FONT_BODY,
        font_size=12,
        bold=True,
        color=BLACK,
        margins=(0, 0, 0, 0),
    )
    chart_top = top + Inches(0.34)
    row_gap = Inches(0.14)
    row_height = Inches(0.42)
    label_width = Inches(1.95)
    rail_width = width - label_width - Inches(0.70)

    for index, (label, value, color) in enumerate(items):
        y = chart_top + index * (row_height + row_gap)
        add_textbox(
            slide,
            left,
            y + Inches(0.08),
            label_width - Inches(0.10),
            Inches(0.16),
            label,
            font_name=FONT_BODY,
            font_size=11,
            color=BLACK,
            margins=(0, 0, 0, 0),
        )

        rail = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            emu(left + label_width),
            emu(y + Inches(0.10)),
            emu(rail_width),
            emu(Inches(0.18)),
        )
        fill_shape(rail, LIGHT_GRAY)
        line_shape(rail, LIGHT_GRAY, 0.4)
        rail.adjustments[0] = 0.2

        bar_width = max(Inches(0.08), rail_width * value / max_value)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            emu(left + label_width),
            emu(y + Inches(0.10)),
            emu(bar_width),
            emu(Inches(0.18)),
        )
        fill_shape(bar, color)
        line_shape(bar, color, 0.4)
        bar.adjustments[0] = 0.2

        add_textbox(
            slide,
            left + label_width + rail_width + Inches(0.05),
            y + Inches(0.04),
            Inches(0.62),
            Inches(0.18),
            f"{value:.3f} с",
            font_name=FONT_BODY,
            font_size=10.5,
            bold=True,
            color=BLACK,
            margins=(0, 0, 0, 0),
        )


def parse_metadata() -> dict[str, str]:
    content = METADATA_FILE.read_text(encoding="utf-8")
    result = {}
    for key in [
        "UniversityName",
        "InstituteName",
        "DepartmentName",
        "WorkTitle",
        "StudentName",
        "StudentGroup",
        "SupervisorName",
        "CityName",
        "WorkYear",
    ]:
        match = re.search(rf"\\newcommand{{\\{key}}}{{(.+?)}}", content)
        if not match:
            raise ValueError(f"Не найден ключ {key}")
        result[key] = match.group(1)
    return result


def load_results() -> dict:
    with RESULTS_FILE.open(encoding="utf-8") as file:
        return json.load(file)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    meta = parse_metadata()
    results = load_results()

    slide_title(prs, meta)
    slide_relevance(prs)
    slide_goal(prs)
    slide_decisions(prs)
    slide_architecture(prs)
    slide_pipeline(prs)
    slide_implementation(prs)
    slide_stack(prs)
    slide_methodology(prs)
    slide_results(prs, results)
    slide_conclusion(prs)
    return prs


def slide_title(prs: Presentation, meta: dict[str, str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_reference_badge(slide, large=True)

    add_textbox(
        slide,
        Inches(1.18),
        Inches(0.22),
        Inches(8.05),
        Inches(0.18),
        meta["UniversityName"],
        font_name=FONT_BODY,
        font_size=10,
        color=BLACK,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )
    add_textbox(
        slide,
        Inches(1.45),
        Inches(0.42),
        Inches(7.55),
        Inches(0.18),
        meta["InstituteName"],
        font_name=FONT_BODY,
        font_size=12,
        color=BLACK,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )
    add_textbox(
        slide,
        Inches(1.20),
        Inches(0.62),
        Inches(8.00),
        Inches(0.18),
        meta["DepartmentName"],
        font_name=FONT_BODY,
        font_size=11,
        color=BLACK,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )

    add_textbox(
        slide,
        Inches(1.25),
        Inches(1.72),
        Inches(7.55),
        Inches(0.28),
        "Магистерская работа на тему:",
        font_name=FONT_BODY,
        font_size=18,
        color=BLACK,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )

    title = (
        "«Разработка и экспериментальная оценка платформы\n"
        "видеоаналитики реального времени для обнаружения,\n"
        "сопровождения и событийного анализа объектов\n"
        "в системах охранного видеонаблюдения»"
    )
    add_textbox(
        slide,
        Inches(1.08),
        Inches(2.06),
        Inches(7.90),
        Inches(1.18),
        title,
        font_name=FONT_TITLE,
        font_size=20,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
        margins=(0, 0, 0, 0),
    )

    add_label_value_line(
        slide,
        Inches(0.72),
        Inches(4.18),
        Inches(8.10),
        "Студент группы " + meta["StudentGroup"] + ": ",
        meta["StudentName"],
        font_size=13,
    )
    add_label_value_line(
        slide,
        Inches(0.72),
        Inches(4.44),
        Inches(8.10),
        "Научный руководитель: ",
        meta["SupervisorName"],
        font_size=13,
    )

    add_textbox(
        slide,
        Inches(4.05),
        Inches(5.00),
        Inches(1.80),
        Inches(0.20),
        f"{meta['CityName']} – {meta['WorkYear']}",
        font_name=FONT_BODY,
        font_size=12,
        color=BLACK,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )


def slide_relevance(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 2, "Актуальность работы")

    add_content_box(slide, Inches(0.62), Inches(1.05), Inches(4.20), Inches(3.60))
    add_textbox(
        slide,
        Inches(0.80),
        Inches(1.20),
        Inches(3.75),
        Inches(0.20),
        "Основные тезисы",
        font_name=FONT_BODY,
        font_size=12,
        bold=True,
        color=BLUE,
        margins=(0, 0, 0, 0),
    )
    bullets = slide.shapes.add_textbox(emu(Inches(0.80)), emu(Inches(1.48)), emu(Inches(3.72)), emu(Inches(2.85)))
    add_bullets(
        bullets,
        [
            "Растет количество камер и объем видеоданных.",
            "Ручной мониторинг не обеспечивает устойчивое наблюдение без пропуска событий.",
            "Для охраняемых зон нужны не только детекция, но и сопровождение, событийный анализ и уведомления.",
            "Особенно сложны сценарии с малозаметными и динамичными объектами, включая БПЛА.",
        ],
        font_size=13,
        color=BLACK,
        space_after=10,
    )

    add_content_box(slide, Inches(5.02), Inches(1.05), Inches(4.35), Inches(3.60))
    add_textbox(
        slide,
        Inches(5.18),
        Inches(1.20),
        Inches(3.95),
        Inches(0.20),
        "Почему ручной мониторинг перестает работать",
        font_name=FONT_BODY,
        font_size=12,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )

    camera_positions = [1.62, 2.24, 2.86]
    for idx, y in enumerate(camera_positions, start=1):
        add_arrow_box(slide, Inches(5.42), Inches(y), Inches(1.10), Inches(0.38), f"Камера {idx}", font_size=11)
        add_connector(slide, Inches(6.54), Inches(y + 0.19), Inches(6.98), Inches(2.50), color=BLUE)

    add_arrow_box(
        slide,
        Inches(6.98),
        Inches(2.20),
        Inches(1.45),
        Inches(0.58),
        "Оператор",
        fill=WHITE,
        border=BLUE,
        font_size=13,
    )
    add_connector(slide, Inches(8.43), Inches(2.49), Inches(8.80), Inches(2.49), color=BLUE)
    add_arrow_box(
        slide,
        Inches(8.82),
        Inches(2.18),
        Inches(0.42),
        Inches(0.42),
        "!",
        fill=LIGHT_BLUE,
        border=BLUE,
        font_size=16,
    )
    add_textbox(
        slide,
        Inches(7.00),
        Inches(3.12),
        Inches(2.05),
        Inches(0.44),
        "Перегрузка, рост времени реакции\nи риск пропуска событий",
        font_name=FONT_BODY,
        font_size=11,
        color=GRAY,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )

    add_note(
        slide,
        Inches(0.62),
        Inches(4.82),
        Inches(8.75),
        "Вывод: актуальна не отдельная модель детекции, а платформа реального времени с устойчивым событийным контуром.",
    )


def slide_goal(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 3, "Цель и задачи исследования")

    add_content_box(slide, Inches(0.62), Inches(1.05), Inches(4.15), Inches(3.80))
    add_textbox(
        slide,
        Inches(0.78),
        Inches(1.20),
        Inches(3.75),
        Inches(0.18),
        "Цель, объект и предмет",
        font_name=FONT_BODY,
        font_size=12,
        bold=True,
        color=BLUE,
        margins=(0, 0, 0, 0),
    )
    goal_box = slide.shapes.add_textbox(emu(Inches(0.78)), emu(Inches(1.52)), emu(Inches(3.75)), emu(Inches(3.00)))
    add_paragraph_lines(
        goal_box,
        [
            "Цель: повышение эффективности автоматизированного мониторинга охраняемых зон за счет разработки и экспериментальной оценки платформы видеоаналитики реального времени.",
            "Объект: процессы автоматизированного мониторинга видеопотоков в системах охранного видеонаблюдения.",
            "Предмет: методы, модели и архитектурные решения построения платформы.",
        ],
        font_size=12.5,
        color=BLACK,
        space_after=12,
    )

    add_content_box(slide, Inches(5.02), Inches(1.05), Inches(4.35), Inches(3.80))
    add_textbox(
        slide,
        Inches(5.18),
        Inches(1.20),
        Inches(3.95),
        Inches(0.18),
        "Задачи",
        font_name=FONT_BODY,
        font_size=12,
        bold=True,
        color=BLUE,
        margins=(0, 0, 0, 0),
    )
    tasks_box = slide.shapes.add_textbox(emu(Inches(5.18)), emu(Inches(1.48)), emu(Inches(3.92)), emu(Inches(3.05)))
    add_bullets(
        tasks_box,
        [
            "Сформировать требования к платформе.",
            "Обосновать выбор методов детекции, трекинга и событийного анализа.",
            "Разработать архитектуру и алгоритмический конвейер.",
            "Реализовать программную платформу.",
            "Провести экспериментальную оценку.",
            "Определить направления дальнейшего развития.",
        ],
        font_size=12.5,
        color=BLACK,
        space_after=8,
    )


def slide_decisions(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 4, "Обоснование выбранных решений")

    data = [
        ["Контур", "Выбранное решение", "Почему выбрано"],
        ["Детекция", "YOLOv8n", "Баланс между скоростью и качеством для потоковой обработки видео."],
        ["Трекинг", "DeepSORT", "Сохраняет идентичность объектов между кадрами."],
        ["События", "Правиловой подход", "Интерпретируемые правила для охранного видеонаблюдения."],
        ["Интеграция", "RabbitMQ + transactional outbox", "Исключает рассогласование между записью события и публикацией."],
        ["Данные", "PostgreSQL + файловый слой", "Структурированные данные хранятся в БД, тяжелые артефакты вынесены отдельно."],
    ]
    add_table(slide, Inches(0.55), Inches(1.05), Inches(8.90), Inches(3.25), data, [1.35, 2.15, 5.20])
    add_note(
        slide,
        Inches(0.55),
        Inches(4.50),
        Inches(8.90),
        "Научная новизна: предложен и апробирован сквозной контур capture -> detect -> track -> rules -> outbox -> RabbitMQ/webhook с latest-frame, dedup_key, retry и DLQ.",
    )


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 5, "Архитектура разработанной платформы")

    add_textbox(slide, Inches(0.78), Inches(1.00), Inches(1.50), Inches(0.16), "Control plane", font_size=11, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(0.78), Inches(2.35), Inches(1.50), Inches(0.16), "Data plane", font_size=11, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    add_textbox(slide, Inches(0.78), Inches(4.10), Inches(1.50), Inches(0.16), "Внешний контур", font_size=11, bold=True, color=BLUE, margins=(0, 0, 0, 0))

    add_arrow_box(slide, Inches(0.85), Inches(1.25), Inches(1.50), Inches(0.52), "UI", fill=WHITE, border=BLUE, font_size=13)
    add_arrow_box(slide, Inches(2.75), Inches(1.25), Inches(2.10), Inches(0.52), "API", fill=WHITE, border=BLUE, font_size=13)
    add_arrow_box(slide, Inches(7.20), Inches(1.25), Inches(2.00), Inches(0.52), "Prometheus / Grafana", fill=WHITE, border=BLUE, font_size=12)
    add_connector(slide, Inches(2.35), Inches(1.51), Inches(2.75), Inches(1.51))

    add_arrow_box(slide, Inches(0.85), Inches(2.62), Inches(1.35), Inches(0.58), "Камеры\nRTSP / MP4", fill=LIGHTER_BLUE, border=BLUE, font_size=12)
    add_arrow_box(slide, Inches(2.55), Inches(2.62), Inches(2.05), Inches(0.58), "analytics-worker", fill=LIGHT_BLUE, border=BLUE, font_size=13)
    add_arrow_box(slide, Inches(4.95), Inches(2.62), Inches(1.95), Inches(0.58), "PostgreSQL\n events + outbox", fill=WHITE, border=BLUE, font_size=11)
    add_arrow_box(slide, Inches(7.20), Inches(2.62), Inches(1.80), Inches(0.58), "integration-worker", fill=WHITE, border=BLUE, font_size=12)
    add_arrow_box(slide, Inches(9.15), Inches(2.62), Inches(0.45), Inches(0.58), "MQ", fill=LIGHTER_BLUE, border=BLUE, font_size=11)

    add_connector(slide, Inches(2.20), Inches(2.91), Inches(2.55), Inches(2.91))
    add_connector(slide, Inches(4.60), Inches(2.91), Inches(4.95), Inches(2.91))
    add_connector(slide, Inches(6.90), Inches(2.91), Inches(7.20), Inches(2.91))
    add_connector(slide, Inches(9.00), Inches(2.91), Inches(9.15), Inches(2.91))
    add_connector(slide, Inches(3.65), Inches(1.77), Inches(3.65), Inches(2.62))
    add_connector(slide, Inches(8.20), Inches(1.77), Inches(8.20), Inches(2.62))

    add_arrow_box(slide, Inches(6.92), Inches(4.25), Inches(2.40), Inches(0.58), "Webhook / внешние системы", fill=WHITE, border=BLUE, font_size=12)
    add_connector(slide, Inches(9.38), Inches(3.20), Inches(8.15), Inches(4.25))

    add_note(
        slide,
        Inches(0.85),
        Inches(4.98),
        Inches(5.55),
        "Ключевой элемент архитектуры: события в одной транзакции пишутся в events и event_outbox.",
    )


def slide_pipeline(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 6, "Алгоритмический конвейер платформы")

    steps = [
        "capture",
        "latest-frame",
        "detect",
        "track",
        "rules",
        "outbox",
        "RabbitMQ / webhook",
    ]
    x = 0.58
    for index, step in enumerate(steps):
        add_arrow_box(slide, Inches(x), Inches(1.35), Inches(1.10), Inches(0.58), step, fill=WHITE, border=BLUE, font_size=11)
        if index < len(steps) - 1:
            add_connector(slide, Inches(x + 1.10), Inches(1.64), Inches(x + 1.28), Inches(1.64))
        x += 1.28

    add_note(
        slide,
        Inches(0.58),
        Inches(2.18),
        Inches(8.82),
        "Режим latest-frame ограничивает накопление задержки, а dedup_key подавляет повторную генерацию одинаковых событий.",
    )

    data = [
        ["Правило", "Условие", "Практический смысл"],
        ["zone_enter", "Вход объекта в полигон контроля", "Фиксация попадания в охраняемую зону."],
        ["line_cross", "Пересечение виртуального рубежа", "Контроль направления движения и нарушения границы."],
        ["dwell_time", "Превышение порога пребывания", "Фиксация подозрительной задержки в зоне."],
    ]
    add_table(slide, Inches(0.58), Inches(2.90), Inches(8.82), Inches(1.78), data, [1.50, 3.15, 4.17])


def slide_implementation(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 7, "Реализация и функциональные возможности")

    add_content_box(slide, Inches(0.58), Inches(1.02), Inches(3.20), Inches(3.85))
    add_textbox(slide, Inches(0.78), Inches(1.18), Inches(2.75), Inches(0.18), "Что реализовано", font_size=12, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    bullets = slide.shapes.add_textbox(emu(Inches(0.78)), emu(Inches(1.48)), emu(Inches(2.78)), emu(Inches(3.00)))
    add_bullets(
        bullets,
        [
            "API с JWT-аутентификацией и RBAC.",
            "Управление камерами, зонами, правилами и журналом событий.",
            "Health-check, readiness-check и метрики.",
            "Операторский интерфейс на React + Vite + TypeScript.",
            "Воспроизводимое развёртывание через Docker Compose.",
        ],
        font_size=12.3,
        color=BLACK,
        space_after=8,
    )

    add_image(slide, ASSETS_DIR / "ui-dashboard.png", Inches(4.05), Inches(1.08), Inches(2.55), Inches(1.60), "Dashboard")
    add_image(slide, ASSETS_DIR / "ui-events.png", Inches(6.78), Inches(1.08), Inches(2.55), Inches(1.60), "Events")
    add_image(slide, ASSETS_DIR / "ui-status.png", Inches(5.42), Inches(2.92), Inches(2.55), Inches(1.60), "Status")


def slide_stack(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 8, "Стек технологий")

    data = [
        ["Слой", "Технологии"],
        ["Компьютерное зрение", "YOLOv8n, DeepSORT, PyAV, FFmpeg"],
        ["Backend", "FastAPI, SQLAlchemy, Alembic, Pydantic"],
        ["Интеграция и данные", "RabbitMQ, PostgreSQL, transactional outbox, webhook"],
        ["Frontend", "React, Vite, TypeScript"],
        ["Инфраструктура и мониторинг", "Docker Compose, Prometheus, Grafana"],
    ]
    add_table(slide, Inches(0.68), Inches(1.15), Inches(8.65), Inches(3.45), data, [2.45, 6.20])
    add_note(
        slide,
        Inches(0.68),
        Inches(4.78),
        Inches(8.65),
        "Стек покрывает не только CV-задачи, но и эксплуатационный контур: безопасность, интеграцию, мониторинг и воспроизводимое развёртывание.",
    )


def slide_methodology(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 9, "Программа апробации и метрики")

    data_profiles = [
        ["Профиль", "Назначение"],
        ["CPU dev, 1–3 камеры", "Функциональная и регрессионная проверка, серия прогонов >= 30 минут."],
        ["Многокамерная серия", "Перспективная проверка p95 object->event, очередей и dropped frames."],
    ]
    add_table(slide, Inches(0.55), Inches(1.05), Inches(8.90), Inches(1.35), data_profiles, [2.25, 6.65])

    add_content_box(slide, Inches(0.55), Inches(2.62), Inches(4.15), Inches(2.10))
    add_textbox(slide, Inches(0.74), Inches(2.78), Inches(3.75), Inches(0.16), "Проверяемые сценарии", font_size=12, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    scenarios = slide.shapes.add_textbox(emu(Inches(0.74)), emu(Inches(3.05)), emu(Inches(3.72)), emu(Inches(1.40)))
    add_bullets(
        scenarios,
        [
            "E2E для zone_enter, line_cross, dwell_time.",
            "Отказ RabbitMQ и восстановление публикации.",
            "Отказ webhook и повторная доставка.",
            "Проверка JWT / RBAC.",
        ],
        font_size=11.5,
        color=BLACK,
        space_after=7,
    )

    data_metrics = [
        ["Метрика / критерий", "Целевое значение"],
        ["CPU p95 object->event", "<= 2.0 с"],
        ["Многокамерный p95 object->event", "следующая серия"],
        ["DLQ", "0 в штатном режиме"],
        ["viewer_create_camera_status", "403"],
        ["Dropped frames", "Без неограниченного роста"],
    ]
    add_table(slide, Inches(4.95), Inches(2.62), Inches(4.50), Inches(2.10), data_metrics, [2.65, 1.85])


def slide_results(prs: Presentation, results: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 10, "Основные результаты апробации")

    summary = (
        f"Acceptance-прогон от 8 марта 2026 года завершён со статусом {results['result']} "
        f"при общей длительности {results['duration_sec']:.3f} с."
    )
    add_note(slide, Inches(0.55), Inches(1.00), Inches(8.90), summary)

    add_content_box(slide, Inches(0.55), Inches(1.72), Inches(4.85), Inches(3.15))
    add_horizontal_bar_chart(
        slide,
        Inches(0.78),
        Inches(1.92),
        Inches(4.35),
        Inches(2.55),
        [
            ("Functional delivery", results["timings_sec"]["functional_delivery"], BLUE),
            ("Webhook recovery", results["timings_sec"]["webhook_recovery"], ORANGE),
            ("RabbitMQ recovery", results["timings_sec"]["rabbitmq_recovery"], RED),
        ],
        max_value=10.0,
        title="Время доставки и восстановления",
    )
    add_textbox(
        slide,
        Inches(0.78),
        Inches(4.45),
        Inches(4.20),
        Inches(0.16),
        "Чем меньше время, тем устойчивее работает полный событийный контур.",
        font_name=FONT_BODY,
        font_size=10,
        color=GRAY,
        margins=(0, 0, 0, 0),
    )

    data = [
        ["Проверка", "Фактический результат"],
        ["Security", f"viewer_create_camera_status = {results['checks']['security']['viewer_create_camera_status']}"],
        ["Functional E2E", "outbox = published, retry_count = 0, success = 1, failed = 0"],
        ["RabbitMQ outage", "Во время отказа status = retry, после восстановления status = published"],
        ["Webhook retry", "failed_attempts_before_recovery = 1, success = 1, failed = 2"],
    ]
    add_table(slide, Inches(5.62), Inches(1.72), Inches(3.83), Inches(2.50), data, [1.35, 2.48])
    add_note(
        slide,
        Inches(5.62),
        Inches(4.30),
        Inches(3.83),
        "Ограничение текущего этапа: длительные многокамерные прогоны не заявляются как фактический результат.",
        height=0.70,
        font_size=10,
    )


def slide_conclusion(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 11, "Заключение")

    add_content_box(slide, Inches(0.58), Inches(1.02), Inches(4.35), Inches(3.62))
    add_textbox(slide, Inches(0.78), Inches(1.18), Inches(3.95), Inches(0.18), "Итоги работы", font_size=12, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    bullets = slide.shapes.add_textbox(emu(Inches(0.78)), emu(Inches(1.46)), emu(Inches(3.92)), emu(Inches(2.85)))
    add_bullets(
        bullets,
        [
            "Разработана платформа видеоаналитики реального времени для мониторинга охраняемых зон.",
            "Обоснован стек методов и реализован сквозной конвейер обработки.",
            "Реализованы механизмы надежной доставки, безопасности и наблюдаемости.",
            "Подготовлена и частично проведена экспериментальная оценка.",
        ],
        font_size=12.5,
        color=BLACK,
        space_after=10,
    )

    add_content_box(slide, Inches(5.15), Inches(1.02), Inches(4.20), Inches(1.58))
    add_textbox(slide, Inches(5.35), Inches(1.18), Inches(3.82), Inches(0.18), "Научная новизна", font_size=12, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    novelty = slide.shapes.add_textbox(emu(Inches(5.35)), emu(Inches(1.48)), emu(Inches(3.75)), emu(Inches(0.92)))
    add_paragraph_lines(
        novelty,
        [
            "Новизна состоит в согласовании real-time CV-конвейера с надежной событийной доставкой.",
            "Качество платформы оценивается по latency, надежности и RBAC, а не только по CV-метрикам.",
        ],
        font_size=10.5,
        color=BLACK,
        space_after=4,
    )

    add_content_box(slide, Inches(5.15), Inches(2.78), Inches(4.20), Inches(2.12))
    add_textbox(slide, Inches(5.35), Inches(2.88), Inches(3.82), Inches(0.18), "Дальнейшее развитие", font_size=12, bold=True, color=BLUE, margins=(0, 0, 0, 0))
    future = slide.shapes.add_textbox(emu(Inches(5.35)), emu(Inches(3.16)), emu(Inches(3.75)), emu(Inches(1.50)))
    add_bullets(
        future,
        [
            "Дообучение детектора на предметных датасетах.",
            "Межкамерная re-identification.",
            "Ускорение инференса через ONNX Runtime или TensorRT.",
            "Расширение масштабируемости инфраструктуры.",
        ],
        font_size=10.4,
        color=BLACK,
        space_after=4,
    )

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(Inches(0.58)), emu(Inches(5.02)), emu(Inches(8.77)), emu(Inches(0.34)))
    fill_shape(band, BLUE)
    line_shape(band, BLUE, 0.5)
    add_textbox(
        slide,
        Inches(0.58),
        Inches(5.10),
        Inches(8.77),
        Inches(0.12),
        "Поставленная цель магистерской работы достигнута.",
        font_name=FONT_BODY,
        font_size=12,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        margins=(0, 0, 0, 0),
    )


if __name__ == "__main__":
    presentation = build_presentation()
    presentation.save(OUTPUT_FILE)
    print(OUTPUT_FILE)
